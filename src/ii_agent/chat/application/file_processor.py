"""Unified file processor for routing and content extraction."""

import io
from typing import List, Set, Optional, BinaryIO, Tuple
from dataclasses import dataclass, field
from abc import ABC, abstractmethod
from pathlib import Path
from urllib.parse import urlparse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from ii_agent.files.models import FileAsset
from ii_agent.chat.types import BinaryContent, TextContent
from ii_agent.core.storage.client import get_storage
from ii_agent.core.config.llm_config import APITypes
from ii_agent.core.logger import logger

# ============================================================================
# Exceptions
# ============================================================================


class ImageCompressionError(Exception):
    """Raised when image compression fails to meet provider size limits."""


# ============================================================================
# Provider Image Limits
# ============================================================================

# Provider-specific limits (raw bytes before base64 encoding)
# Anthropic: 5MB base64 limit → ~3.75MB raw
# OpenAI/Gemini: 10MB limit
PROVIDER_IMAGE_LIMITS = {
    APITypes.ANTHROPIC: int(5 * 1024 * 1024 * 3 / 4),  # ~3.75MB (5MB base64)
    APITypes.OPENAI: 10 * 1024 * 1024,  # 10MB
    APITypes.GEMINI: 10 * 1024 * 1024,  # 10MB
    APITypes.CUSTOM: 10 * 1024 * 1024,  # 10MB default
}
DEFAULT_IMAGE_LIMIT = int(5 * 1024 * 1024 * 3 / 4)  # Conservative default (~3.75MB)


# File size thresholds
SIZE_THRESHOLD_MB = 50

# Token thresholds for content routing
# Leave ~28k tokens for conversation context and system prompt
# Claude has 128k context, so we use 100k as the max for file content
MAX_CONTENT_TOKENS = 100000

# Approximate tokens per character (conservative estimate: 1 token ≈ 3 chars)
CHARS_PER_TOKEN = 3

# PDF page limit for Anthropic Claude API
MAX_PDF_PAGES = 100

# ============================================================================
# Data Models
# ============================================================================


@dataclass
class ProcessedFiles:
    """Result of file processing.

    Contains content parts ready to add to message and file IDs for different strategies.
    """

    binary_parts: List[BinaryContent]  # PDF/images as BinaryContent (for messages)
    text_parts: List[TextContent]  # Text files content (for messages)
    large_file_ids: Set[str]  # Large files for vector store
    large_file_info: List[dict]  # Info about large files for logging
    skipped_files: List[dict]  # Files that couldn't be processed
    total_text_tokens: int = 0  # Total estimated tokens from text content
    vector_store_reason: List[dict] = field(
        default_factory=list
    )  # Reasons files were routed to vector store


# ============================================================================
# Helper Functions
# ============================================================================


def estimate_tokens(text: str) -> int:
    """
    Estimate the number of tokens in text.

    Uses a conservative estimate of 1 token per 3 characters.
    This is a rough approximation but sufficient for routing decisions.
    Uses ceiling division to avoid underestimation.
    """
    import math

    return math.ceil(len(text) / CHARS_PER_TOKEN)


def get_pdf_page_count(file_bytes: bytes) -> int:
    """
    Get the number of pages in a PDF file.

    Args:
        file_bytes: PDF file content as bytes

    Returns:
        Number of pages, or -1 if unable to determine
    """
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            page_count = len(doc)
            return page_count
        finally:
            doc.close()
    except ImportError:
        logger.warning("[PDF] PyMuPDF not available, cannot count PDF pages")
        return -1
    except Exception as e:
        logger.error(f"[PDF] Error counting PDF pages: {e}")
        return -1


def extract_pdf_text(file_bytes: bytes) -> Optional[str]:
    """
    Extract text from a PDF file.

    Args:
        file_bytes: PDF file content as bytes

    Returns:
        Extracted text or None if extraction fails
    """
    try:
        import fitz  # PyMuPDF

        text_content = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        try:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    text_content.append(f"[Page {page_num + 1}]\n{text}")

            return "\n\n".join(text_content) if text_content else None
        finally:
            doc.close()

    except ImportError:
        logger.warning("[PDF] PyMuPDF not available, cannot extract PDF text")
        return None
    except Exception as e:
        logger.error(f"[PDF] Error extracting PDF text: {e}")
        return None


def is_text_extractable(content_type: Optional[str], file_name: str) -> bool:
    """
    Check if file supports text extraction.

    Uses ContentExtractorFactory to determine if an extractor exists.
    """
    extractor = ContentExtractorFactory.get_extractor(content_type, file_name)
    return extractor is not None


def is_binary_file(content_type: Optional[str], file_name: str) -> bool:
    """Check if file is PDF or image (supported for BinaryContent)."""
    _binary_extensions = (".pdf", ".png", ".jpg", ".jpeg", ".gif", ".webp", ".heic", ".heif")
    if not content_type or content_type == "application/octet-stream":
        file_lower = file_name.lower()
        return file_lower.endswith(_binary_extensions)

    # PDF
    if content_type == "application/pdf":
        return True

    # Images
    if content_type.startswith("image/"):
        return True

    return False


def is_remote_url(path: str) -> bool:
    """Check if the given path is an HTTP(S) URL."""
    parsed = urlparse(path)
    return parsed.scheme in {"http", "https"}


def compress_image_for_provider(
    file_bytes: bytes,
    mime_type: str,
    target_limit: int,
) -> Tuple[bytes, str]:
    """
    Compress image to meet provider-specific size limits.

    Uses progressive compression strategy:
    1. Check if already under limit → return unchanged
    2. Try different max dimensions (4K → 1K)
    3. Try different quality levels (95 → 35)
    4. If all attempts fail → raise ImageCompressionError

    Args:
        file_bytes: Original image bytes
        mime_type: Original MIME type
        target_limit: Target size limit in bytes

    Returns:
        Tuple of (compressed_bytes, mime_type)

    Raises:
        ImageCompressionError: If compression cannot meet the limit
    """
    from PIL import Image

    # Detect HEIC via magic bytes when MIME type doesn't indicate it
    # (e.g. application/octet-stream or missing MIME type)
    from ii_agent.agents.utils.heic import is_heic_format

    is_heic = is_heic_format(mime_type=mime_type, image_bytes=file_bytes)
    if is_heic and mime_type not in ("image/heic", "image/heif"):
        mime_type = "image/heic"

    # Register HEIC/HEIF support via pillow-heif plugin
    if is_heic:
        try:
            from pillow_heif import register_heif_opener

            register_heif_opener()
        except ImportError:
            raise ImageCompressionError(
                "HEIC/HEIF support requires the pillow-heif package. "
                "Please install it with: pip install pillow-heif"
            )

    # Check if already under limit (skip for HEIC since providers don't support it)
    if len(file_bytes) <= target_limit and not is_heic:
        logger.debug(
            f"[IMAGE_COMPRESSION] Image already under limit "
            f"({len(file_bytes)} bytes <= {target_limit} bytes)"
        )
        return file_bytes, mime_type

    original_size = len(file_bytes)
    logger.info(
        f"[IMAGE_COMPRESSION] Starting compression: {original_size} bytes → target {target_limit} bytes"
    )

    try:
        from PIL import ImageOps

        # Open the image
        img = Image.open(io.BytesIO(file_bytes))

        # Apply EXIF orientation (HEIC photos from iPhones carry rotation metadata)
        img = ImageOps.exif_transpose(img)

        # Convert RGBA/P to RGB (for JPEG compression)
        if img.mode in ("RGBA", "P"):
            logger.debug(f"[IMAGE_COMPRESSION] Converting {img.mode} to RGB")
            img = img.convert("RGB")

        original_dimensions = img.size

        # Progressive compression strategy
        max_dimensions = [4096, 3072, 2048, 1536, 1024]
        quality_levels = [95, 85, 75, 65, 55, 45, 35]

        for max_dim in max_dimensions:
            # Resize if needed (maintaining aspect ratio)
            resized_img = img.copy()
            if max(img.size) > max_dim:
                resized_img.thumbnail((max_dim, max_dim), Image.Resampling.LANCZOS)
                logger.debug(
                    f"[IMAGE_COMPRESSION] Resized from {original_dimensions} to {resized_img.size}"
                )

            for quality in quality_levels:
                # Compress as JPEG
                buffer = io.BytesIO()
                resized_img.save(buffer, format="JPEG", quality=quality, optimize=True)
                compressed_bytes = buffer.getvalue()

                if len(compressed_bytes) <= target_limit:
                    logger.info(
                        f"[IMAGE_COMPRESSION] Success: {original_size} → {len(compressed_bytes)} bytes "
                        f"(max_dim={max_dim}, quality={quality}, dimensions={resized_img.size})"
                    )
                    return compressed_bytes, "image/jpeg"

                logger.debug(
                    f"[IMAGE_COMPRESSION] Attempt: max_dim={max_dim}, quality={quality} → "
                    f"{len(compressed_bytes)} bytes (still too large)"
                )

        # All attempts failed
        raise ImageCompressionError(
            f"Image could not be compressed to meet the {target_limit / 1024 / 1024:.1f}MB limit. "
            f"Original size: {original_size / 1024 / 1024:.1f}MB. "
            f"Please upload a smaller image."
        )

    except ImageCompressionError:
        raise
    except Exception as e:
        logger.opt(exception=True).error(f"[IMAGE_COMPRESSION] Error during compression: {e}")
        raise ImageCompressionError(
            f"Failed to compress image: {str(e)}. Please upload a smaller image."
        )


async def process_files_for_message(
    db_session: AsyncSession,
    file_ids: List[str],
    storage,
    session_id: str,
    api_type: Optional[APITypes] = None,
) -> ProcessedFiles:
    """
    Process files and prepare content parts for message.

    Strategy:
    - Files > 50MB → Large files → Vector store + FileSearchTool
    - PDFs with >100 pages → Extract text → Check token limits
    - PDF/images ≤ 50MB with ≤100 pages → BinaryContent (base64 encoded)
    - Text files ≤ 50MB → TextContent (extract text and add to message)
    - If total extracted text > 100k tokens → Route excess to vector store
    - Other files ≤ 50MB → Skip (unsupported)

    Args:
        db_session: Database session
        file_ids: List of file IDs to process
        storage: Storage client for reading files
        session_id: Session ID for logging
        api_type: Optional API type for provider-specific image compression

    Returns:
        ProcessedFiles with content parts and routing info

    Raises:
        ImageCompressionError: If an image cannot be compressed to meet provider limits
    """
    if not file_ids:
        logger.info(f"[FILE_PROCESSOR] session={session_id}: No files to process")
        return ProcessedFiles(
            binary_parts=[],
            text_parts=[],
            large_file_ids=set(),
            large_file_info=[],
            skipped_files=[],
        )

    logger.info(f"[FILE_PROCESSOR] session={session_id}: Processing {len(file_ids)} files")

    # Load all files in one query
    stmt = select(FileAsset).where(FileAsset.id.in_(file_ids))
    result = await db_session.execute(stmt)
    file_uploads = result.scalars().all()

    binary_parts = []
    text_parts = []
    large_file_ids = set()
    large_file_info = []
    skipped_files = []
    vector_store_reason = []
    total_text_tokens = 0

    for file_upload in file_uploads:
        file_size_mb = file_upload.file_size / 1024 / 1024

        # Strategy 1: Large files → Vector store
        if file_size_mb > SIZE_THRESHOLD_MB:
            large_file_ids.add(file_upload.id)
            large_file_info.append(
                {
                    "file_id": file_upload.id,
                    "file_name": file_upload.file_name,
                    "size_kb": f"{file_upload.file_size / 1024:.2f}",
                }
            )
            vector_store_reason.append(
                {
                    "file_name": file_upload.file_name,
                    "reason": f"File too large ({file_size_mb:.1f}MB > {SIZE_THRESHOLD_MB}MB)",
                }
            )
            logger.info(
                f"[FILE_PROCESSOR] LARGE: {file_upload.file_name} "
                f"({file_size_mb:.2f}MB) → Vector store"
            )
            continue

        # Strategy 2: Small PDF/images → BinaryContent (with page limit check for PDFs)
        if is_binary_file(file_upload.content_type, file_upload.file_name):
            try:
                import anyio
                import httpx

                if is_remote_url(file_upload.storage_path):
                    async with httpx.AsyncClient(follow_redirects=True, timeout=30.0) as client:
                        response = await client.get(file_upload.storage_path)
                        response.raise_for_status()
                        file_bytes = response.content
                        mime_type = (
                            file_upload.content_type
                            or response.headers.get("content-type")
                            or "application/octet-stream"
                        )
                else:
                    # All files use unified storage
                    file_content = await anyio.to_thread.run_sync(
                        get_storage().read, file_upload.storage_path
                    )
                    file_bytes = file_content.read()
                    file_content.close()
                    mime_type = file_upload.content_type

                # Correct generic MIME types for images detected by extension or magic bytes
                if not mime_type or mime_type == "application/octet-stream":
                    from ii_agent.agents.utils.heic import is_heic_format as _is_heic

                    if file_bytes and _is_heic(image_bytes=file_bytes):
                        mime_type = "image/heic"
                    elif file_upload.file_name:
                        import mimetypes as _mt

                        guessed = _mt.guess_type(file_upload.file_name)[0]
                        if guessed and guessed.startswith("image/"):
                            mime_type = guessed

                # Compress images if needed for provider limits
                # Note: PDFs are not compressed, only images
                if mime_type and mime_type.startswith("image/"):
                    target_limit = PROVIDER_IMAGE_LIMITS.get(api_type, DEFAULT_IMAGE_LIMIT)
                    # compress_image_for_provider raises ImageCompressionError if it fails
                    # Let it propagate to service.py error handler
                    file_bytes, mime_type = compress_image_for_provider(
                        file_bytes, mime_type, target_limit
                    )

                # Check PDF page count for Anthropic API limit
                is_pdf = mime_type == "application/pdf" or file_upload.file_name.lower().endswith(
                    ".pdf"
                )

                if is_pdf:
                    page_count = get_pdf_page_count(file_bytes)
                    logger.info(f"[FILE_PROCESSOR] PDF {file_upload.file_name}: {page_count} pages")

                    # Treat page count failure (-1) or excessive pages as needing text extraction
                    if page_count == -1 or page_count > MAX_PDF_PAGES:
                        # PDF has too many pages - extract text and check token limits
                        logger.info(
                            f"[FILE_PROCESSOR] PDF {file_upload.file_name} has {page_count} pages "
                            f"(> {MAX_PDF_PAGES}), extracting text instead"
                        )

                        extracted_text = extract_pdf_text(file_bytes)
                        if extracted_text:
                            text_tokens = estimate_tokens(extracted_text)

                            # Check if adding this text would exceed token limit
                            if total_text_tokens + text_tokens > MAX_CONTENT_TOKENS:
                                # Route to vector store
                                large_file_ids.add(file_upload.id)
                                large_file_info.append(
                                    {
                                        "file_id": file_upload.id,
                                        "file_name": file_upload.file_name,
                                        "size_kb": f"{file_upload.file_size / 1024:.2f}",
                                    }
                                )
                                vector_store_reason.append(
                                    {
                                        "file_name": file_upload.file_name,
                                        "reason": f"PDF too long ({page_count} pages, ~{text_tokens:,} tokens) - content exceeds context window",
                                    }
                                )
                                logger.info(
                                    f"[FILE_PROCESSOR] PDF {file_upload.file_name}: "
                                    f"~{text_tokens:,} tokens would exceed limit, routing to vector store"
                                )
                            else:
                                # Add as text content
                                formatted_text = (
                                    f"\n\n--- PDF: {file_upload.file_name} ({page_count} pages) ---\n"
                                    f"{extracted_text}\n"
                                    f"--- End of {file_upload.file_name} ---\n"
                                )
                                text_part = TextContent(text=formatted_text)
                                text_parts.append(text_part)
                                total_text_tokens += text_tokens
                                logger.info(
                                    f"[FILE_PROCESSOR] PDF {file_upload.file_name}: "
                                    f"extracted {len(extracted_text)} chars (~{text_tokens:,} tokens) → TextContent"
                                )
                        else:
                            # Text extraction failed, route to vector store
                            large_file_ids.add(file_upload.id)
                            large_file_info.append(
                                {
                                    "file_id": file_upload.id,
                                    "file_name": file_upload.file_name,
                                    "size_kb": f"{file_upload.file_size / 1024:.2f}",
                                }
                            )
                            vector_store_reason.append(
                                {
                                    "file_name": file_upload.file_name,
                                    "reason": f"PDF has {page_count} pages but text extraction failed",
                                }
                            )
                            logger.warning(
                                f"[FILE_PROCESSOR] PDF {file_upload.file_name}: "
                                f"text extraction failed, routing to vector store"
                            )
                        continue

                # PDF with <= 100 pages or image - use binary content
                binary_part = BinaryContent(
                    path=file_upload.storage_path,
                    mime_type=mime_type,
                    data=file_bytes,
                )
                binary_parts.append(binary_part)
                logged_size_mb = file_size_mb if file_size_mb > 0 else len(file_bytes) / 1024 / 1024
                logger.info(
                    f"[FILE_PROCESSOR] BINARY: {file_upload.file_name} "
                    f"({logged_size_mb:.2f}MB, {len(file_bytes)} bytes) → BinaryContent"
                )
            except ImageCompressionError:
                # Re-raise compression errors to be handled by service.py
                raise
            except Exception as e:
                logger.opt(exception=True).error(
                    f"Failed to load binary file {file_upload.file_name}: {e}"
                )
                skipped_files.append(
                    {
                        "file_name": file_upload.file_name,
                        "reason": f"Load error: {str(e)[:100]}",
                    }
                )
            continue

        # Strategy 3: Small text-extractable files → TextContent (with token limit check)
        if is_text_extractable(file_upload.content_type, file_upload.file_name):
            try:
                import anyio

                # All files use unified storage
                file_content = await anyio.to_thread.run_sync(
                    get_storage().read, file_upload.storage_path
                )

                # Extract text using ContentExtractorFactory
                extracted_text = ContentExtractorFactory.extract_content(
                    file_content, file_upload.content_type, file_upload.file_name
                )
                file_content.close()

                if extracted_text and extracted_text.strip():
                    # Estimate tokens for this content
                    text_tokens = estimate_tokens(extracted_text)

                    # Check if adding this text would exceed token limit
                    if total_text_tokens + text_tokens > MAX_CONTENT_TOKENS:
                        # Route to vector store
                        large_file_ids.add(file_upload.id)
                        large_file_info.append(
                            {
                                "file_id": file_upload.id,
                                "file_name": file_upload.file_name,
                                "size_kb": f"{file_upload.file_size / 1024:.2f}",
                            }
                        )
                        vector_store_reason.append(
                            {
                                "file_name": file_upload.file_name,
                                "reason": f"Content too long (~{text_tokens:,} tokens) - would exceed context window limit",
                            }
                        )
                        logger.info(
                            f"[FILE_PROCESSOR] TEXT {file_upload.file_name}: "
                            f"~{text_tokens:,} tokens would exceed limit "
                            f"(current: {total_text_tokens:,}, max: {MAX_CONTENT_TOKENS:,}), "
                            f"routing to vector store"
                        )
                    else:
                        # Format text content with file name
                        formatted_text = (
                            f"\n\n--- File: {file_upload.file_name} ---\n"
                            f"{extracted_text}\n"
                            f"--- End of {file_upload.file_name} ---\n"
                        )

                        text_part = TextContent(text=formatted_text)
                        text_parts.append(text_part)
                        total_text_tokens += text_tokens
                        logger.info(
                            f"[FILE_PROCESSOR] TEXT: {file_upload.file_name} "
                            f"({file_size_mb:.2f}MB, {len(extracted_text)} chars, ~{text_tokens:,} tokens) → TextContent"
                        )
                else:
                    logger.warning(
                        f"[FILE_PROCESSOR] No text extracted from {file_upload.file_name}"
                    )
                    skipped_files.append(
                        {
                            "file_name": file_upload.file_name,
                            "reason": "No text extracted",
                        }
                    )
            except Exception as e:
                logger.opt(exception=True).error(
                    f"Failed to extract text from {file_upload.file_name}: {e}"
                )
                skipped_files.append(
                    {
                        "file_name": file_upload.file_name,
                        "reason": f"Extraction error: {str(e)[:100]}",
                    }
                )
            continue

        # Strategy 4: Unsupported file types
        logger.warning(
            f"[FILE_PROCESSOR] SKIP: {file_upload.file_name} "
            f"({file_size_mb:.2f}MB) → Unsupported type: {file_upload.content_type}"
        )
        skipped_files.append(
            {
                "file_name": file_upload.file_name,
                "reason": f"Unsupported type: {file_upload.content_type}",
            }
        )

    # Summary log
    logger.info(
        f"[FILE_PROCESSOR] session={session_id}: "
        f"SUMMARY: {len(binary_parts)} binary, {len(text_parts)} text (~{total_text_tokens:,} tokens), "
        f"{len(large_file_ids)} large (vector store), {len(skipped_files)} skipped"
    )

    if skipped_files:
        logger.warning(f"[FILE_PROCESSOR] Skipped files: {[f['file_name'] for f in skipped_files]}")

    if vector_store_reason:
        logger.info(f"[FILE_PROCESSOR] Vector store routing reasons: {vector_store_reason}")

    return ProcessedFiles(
        binary_parts=binary_parts,
        text_parts=text_parts,
        large_file_ids=large_file_ids,
        large_file_info=large_file_info,
        skipped_files=skipped_files,
        total_text_tokens=total_text_tokens,
        vector_store_reason=vector_store_reason,
    )


# ============================================================================
# Content Extractors
# ============================================================================
# Organized by category following OpenAI supported file types:
# 1. Text formats: .txt, .md, .log
# 2. Code formats: .py, .js, .java, .cpp, .c, .cs, .php, .rb, .go, .rs, .ts, .tsx, .jsx
# 3. Document formats: .pdf, .docx, .doc, .pptx, .ppt
# 4. Data formats: .json, .csv, .xml, .xlsx, .xls
# 5. Web formats: .html, .htm
# 6. Rich text: .rtf
# ============================================================================


class ContentExtractor(ABC):
    """Base class for file content extraction."""

    @abstractmethod
    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """
        Extract text content from file object.

        Args:
            file_obj: Binary file object (from storage.read())

        Returns:
            Extracted text content or None if extraction fails
        """
        pass


# ============================================================================
# Category 1: Plain Text Extractors
# ============================================================================


class TextExtractor(ContentExtractor):
    """Extractor for plain text files (.txt, .log, .text)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract text from plain text file."""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            return content.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting text: {e}")
            return None


class MarkdownExtractor(ContentExtractor):
    """Extractor for markdown files (.md, .markdown)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract content from markdown file."""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            return content.decode("utf-8", errors="ignore")
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting markdown: {e}")
            return None


# ============================================================================
# Category 2: Code File Extractors
# ============================================================================


class CodeExtractor(ContentExtractor):
    """Extractor for source code files (various languages)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract code from source file."""
        try:
            file_obj.seek(0)
            content = file_obj.read()
            # Try UTF-8 first, fallback to latin-1
            try:
                return content.decode("utf-8")
            except UnicodeDecodeError:
                return content.decode("latin-1", errors="ignore")
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting code: {e}")
            return None


# ============================================================================
# Category 3: Document Format Extractors
# ============================================================================


# class PDFExtractor(ContentExtractor):
#     """Extractor for PDF files using PyMuPDF (fitz)."""

#     def extract(self, file_obj: BinaryIO) -> Optional[str]:
#         """Extract text from PDF file."""
#         try:
#             import fitz  # PyMuPDF

#             file_obj.seek(0)
#             text_content = []
#             doc = fitz.open(stream=file_obj.read(), filetype="pdf")

#             for page_num, page in enumerate(doc):
#                 # Extract text from page
#                 text = page.get_text()
#                 if text.strip():
#                     text_content.append(text)

#                 # Log if page has images (for debugging)
#                 image_count = len(page.get_images())
#                 if image_count > 0:
#                     logger.debug(
#                         f"[EXTRACTOR] PDF page {page_num + 1} contains {image_count} image(s)"
#                     )

#             doc.close()
#             return "\n\n".join(text_content) if text_content else None

#         except ImportError:
#             logger.warning("[EXTRACTOR] PyMuPDF not available, cannot extract PDF")
#             return None
#         except Exception as e:
#             logger.error(f"[EXTRACTOR] Error extracting PDF: {e}")
#             return None


class WordExtractor(ContentExtractor):
    """Extractor for Word documents (.docx, .doc)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract text from Word document."""
        try:
            from docx import Document

            file_obj.seek(0)
            doc = Document(file_obj)
            text_content = []

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    if any(row_text):
                        text_content.append(" | ".join(row_text))

            return "\n".join(text_content) if text_content else None

        except ImportError:
            logger.warning("[EXTRACTOR] python-docx not available, cannot extract Word document")
            return None
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting Word doc: {e}")
            return None


class PowerPointExtractor(ContentExtractor):
    """Extractor for PowerPoint presentations (.pptx, .ppt)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation

            file_obj.seek(0)
            prs = Presentation(file_obj)
            text_content = []

            for slide_num, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                if slide_texts:
                    text_content.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_texts))

            return "\n\n".join(text_content) if text_content else None

        except ImportError:
            logger.warning("[EXTRACTOR] python-pptx not available, cannot extract PowerPoint")
            return None
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting PowerPoint: {e}")
            return None


# ============================================================================
# Category 4: Data Format Extractors
# ============================================================================


class JSONExtractor(ContentExtractor):
    """Extractor for JSON files with pretty formatting."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract and format JSON content."""
        try:
            import json

            file_obj.seek(0)
            content = file_obj.read()

            # Try to parse and pretty-print JSON
            try:
                data = json.loads(content.decode("utf-8", errors="ignore"))
                return json.dumps(data, indent=2, ensure_ascii=False)
            except json.JSONDecodeError:
                # If JSON is invalid, return raw content
                logger.warning("[EXTRACTOR] Invalid JSON, returning raw content")
                return content.decode("utf-8", errors="ignore")

        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting JSON: {e}")
            return None


class CSVExtractor(ContentExtractor):
    """Extractor for CSV files with formatted output."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract and format CSV content."""
        try:
            import csv
            import io

            file_obj.seek(0)
            content = file_obj.read().decode("utf-8", errors="ignore")

            # Parse CSV and format as readable text
            reader = csv.reader(io.StringIO(content))
            rows = list(reader)

            if not rows:
                return None

            # Format as markdown table if reasonable size
            if len(rows) <= 100 and len(rows[0]) <= 20:
                formatted_rows = []
                for i, row in enumerate(rows):
                    formatted_rows.append(" | ".join(str(cell) for cell in row))
                    # Add separator after header
                    if i == 0:
                        formatted_rows.append(" | ".join("---" for _ in row))

                return "\n".join(formatted_rows)
            else:
                # For large CSV, just return raw content with warning
                logger.info(f"[EXTRACTOR] Large CSV ({len(rows)} rows), returning raw content")
                return content

        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting CSV: {e}")
            return None


class ExcelExtractor(ContentExtractor):
    """Extractor for Excel files (.xlsx, .xls) using openpyxl."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract and format Excel content."""
        try:
            from openpyxl import load_workbook

            file_obj.seek(0)
            wb = load_workbook(file_obj, data_only=True)
            text_content = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content.append(f"\n=== Sheet: {sheet_name} ===\n")

                # Get all rows with values
                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Skip empty rows
                    if any(cell is not None for cell in row):
                        rows.append(row)

                if not rows:
                    text_content.append("(Empty sheet)")
                    continue

                # Format as markdown table if reasonable size
                if len(rows) <= 100 and len(rows[0]) <= 20:
                    for i, row in enumerate(rows):
                        row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        text_content.append(row_text)
                        # Add separator after header
                        if i == 0:
                            text_content.append(" | ".join("---" for _ in row))
                else:
                    # For large sheets, just show row count
                    text_content.append(f"(Large sheet: {len(rows)} rows, {len(rows[0])} columns)")

            return "\n".join(text_content) if text_content else None

        except ImportError:
            logger.warning("[EXTRACTOR] openpyxl not available, cannot extract Excel")
            return None
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting Excel: {e}")
            return None


class XMLExtractor(ContentExtractor):
    """Extractor for XML files with pretty formatting."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract and format XML content."""
        try:
            import xml.etree.ElementTree as ET
            import xml.dom.minidom as minidom

            file_obj.seek(0)
            content = file_obj.read()

            # Try to parse and pretty-print XML
            try:
                root = ET.fromstring(content)
                xml_str = ET.tostring(root, encoding="unicode")
                pretty_xml = minidom.parseString(xml_str).toprettyxml(indent="  ")
                return pretty_xml
            except ET.ParseError:
                # If XML is invalid, return raw content
                logger.warning("[EXTRACTOR] Invalid XML, returning raw content")
                return content.decode("utf-8", errors="ignore")

        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting XML: {e}")
            return None


# ============================================================================
# Category 5: Web Format Extractors
# ============================================================================


class HTMLExtractor(ContentExtractor):
    """Extractor for HTML files with text extraction."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract text content from HTML file."""
        try:
            from bs4 import BeautifulSoup

            file_obj.seek(0)
            content = file_obj.read()
            soup = BeautifulSoup(content, "html.parser")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            return text if text else None

        except ImportError:
            logger.warning("[EXTRACTOR] BeautifulSoup not available, returning raw HTML")
            # Fallback: return raw HTML
            try:
                file_obj.seek(0)
                return file_obj.read().decode("utf-8", errors="ignore")
            except Exception:
                return None
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting HTML: {e}")
            return None


# ============================================================================
# Category 6: Rich Text Format Extractors
# ============================================================================


class RTFExtractor(ContentExtractor):
    """Extractor for Rich Text Format files (.rtf)."""

    def extract(self, file_obj: BinaryIO) -> Optional[str]:
        """Extract text from RTF file."""
        try:
            from striprtf.striprtf import rtf_to_text

            file_obj.seek(0)
            content = file_obj.read().decode("utf-8", errors="ignore")
            text = rtf_to_text(content)

            return text if text.strip() else None

        except ImportError:
            logger.warning("[EXTRACTOR] striprtf not available, cannot extract RTF")
            return None
        except Exception as e:
            logger.error(f"[EXTRACTOR] Error extracting RTF: {e}")
            return None


class ContentExtractorFactory:
    """Factory for creating appropriate content extractor based on file type."""

    # Map MIME types to extractors (organized by category)
    EXTRACTOR_MAP = {
        # Text formats
        "text/plain": TextExtractor,
        "text/markdown": MarkdownExtractor,
        "text/md": MarkdownExtractor,
        # Code formats
        "text/x-python": CodeExtractor,
        "text/x-java": CodeExtractor,
        "text/x-c": CodeExtractor,
        "text/x-c++": CodeExtractor,
        "text/x-csharp": CodeExtractor,
        "text/x-php": CodeExtractor,
        "text/x-ruby": CodeExtractor,
        "text/x-go": CodeExtractor,
        "text/x-rust": CodeExtractor,
        "text/javascript": CodeExtractor,
        "application/javascript": CodeExtractor,
        "application/typescript": CodeExtractor,
        # Document formats
        # "application/pdf": PDFExtractor,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": WordExtractor,
        "application/msword": WordExtractor,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": PowerPointExtractor,
        "application/vnd.ms-powerpoint": PowerPointExtractor,
        # Data formats
        "application/json": JSONExtractor,
        "text/csv": CSVExtractor,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ExcelExtractor,
        "application/vnd.ms-excel": ExcelExtractor,
        "text/xml": XMLExtractor,
        "application/xml": XMLExtractor,
        # Web formats
        "text/html": HTMLExtractor,
        # Rich text
        "application/rtf": RTFExtractor,
        "text/rtf": RTFExtractor,
    }

    @staticmethod
    def get_extractor(content_type: Optional[str], file_name: str) -> Optional[ContentExtractor]:
        """
        Get appropriate content extractor based on content type and filename.

        Args:
            content_type: MIME type
            file_name: Name of the file

        Returns:
            Appropriate ContentExtractor instance or None if no suitable extractor found
        """
        # Try MIME type first
        if content_type and content_type in ContentExtractorFactory.EXTRACTOR_MAP:
            extractor_class = ContentExtractorFactory.EXTRACTOR_MAP[content_type]
            return extractor_class()

        # Fall back to file extension
        suffix = Path(file_name).suffix.lower()

        extension_map = {
            # Text formats
            ".txt": TextExtractor,
            ".text": TextExtractor,
            ".log": TextExtractor,
            ".md": MarkdownExtractor,
            ".markdown": MarkdownExtractor,
            # Code formats
            ".c": CodeExtractor,
            ".cpp": CodeExtractor,
            ".cc": CodeExtractor,
            ".cxx": CodeExtractor,
            ".h": CodeExtractor,
            ".hpp": CodeExtractor,
            ".cs": CodeExtractor,
            ".java": CodeExtractor,
            ".php": CodeExtractor,
            ".rb": CodeExtractor,
            ".go": CodeExtractor,
            ".rs": CodeExtractor,
            ".py": CodeExtractor,
            ".js": CodeExtractor,
            ".ts": CodeExtractor,
            ".tsx": CodeExtractor,
            ".jsx": CodeExtractor,
            ".sh": CodeExtractor,
            ".bash": CodeExtractor,
            ".zsh": CodeExtractor,
            ".sql": CodeExtractor,
            ".swift": CodeExtractor,
            ".kt": CodeExtractor,
            ".scala": CodeExtractor,
            ".r": CodeExtractor,
            ".m": CodeExtractor,
            ".pl": CodeExtractor,
            # Document formats
            # ".pdf": PDFExtractor,
            ".docx": WordExtractor,
            ".doc": WordExtractor,
            ".pptx": PowerPointExtractor,
            ".ppt": PowerPointExtractor,
            # Data formats
            ".json": JSONExtractor,
            ".csv": CSVExtractor,
            ".xlsx": ExcelExtractor,
            ".xls": ExcelExtractor,
            ".xml": XMLExtractor,
            # Web formats
            ".html": HTMLExtractor,
            ".htm": HTMLExtractor,
            # Rich text
            ".rtf": RTFExtractor,
        }

        if suffix in extension_map:
            extractor_class = extension_map[suffix]
            return extractor_class()

        logger.warning(f"[EXTRACTOR] No extractor found for {file_name} (type: {content_type})")
        return None

    @staticmethod
    def extract_content(
        file_obj: BinaryIO, content_type: Optional[str], file_name: str
    ) -> Optional[str]:
        """
        Extract content from file object using appropriate extractor.

        Args:
            file_obj: Binary file object (from storage.read())
            content_type: MIME type
            file_name: Name of the file

        Returns:
            Extracted content or None if extraction fails
        """
        extractor = ContentExtractorFactory.get_extractor(content_type, file_name)

        if not extractor:
            logger.warning(
                f"[EXTRACTOR] Cannot extract content from {file_name}, no suitable extractor found"
            )
            return None

        logger.info(
            f"[EXTRACTOR] Extracting content from {file_name} using {extractor.__class__.__name__}"
        )
        content = extractor.extract(file_obj)

        if content:
            logger.info(
                f"[EXTRACTOR] Successfully extracted {len(content)} characters from {file_name}"
            )
        else:
            logger.warning(f"[EXTRACTOR] No content extracted from {file_name}")

        return content
